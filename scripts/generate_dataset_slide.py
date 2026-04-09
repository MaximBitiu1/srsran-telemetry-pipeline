"""
Standalone generator for the Anomaly Dataset — Labels and Features slide.

Usage:  python3 scripts/generate_dataset_slide.py
Output: docs/dataset_labels_features_slide.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = os.path.dirname(os.path.abspath(__file__)) + "/.."
OUT  = ROOT + "/docs/dataset_labels_features_slide.pptx"

NAVY   = RGBColor(0x1a, 0x1a, 0x3a)
BLUE   = RGBColor(0x1a, 0x6f, 0xaf)
GREEN  = RGBColor(0x2e, 0x8b, 0x57)
PURPLE = RGBColor(0x6a, 0x0d, 0xad)
ORANGE = RGBColor(0xe6, 0x5c, 0x00)
RED    = RGBColor(0xb2, 0x22, 0x22)
WHITE  = RGBColor(0xff, 0xff, 0xff)
LGREY  = RGBColor(0xcc, 0xcc, 0xdd)
YELLOW = RGBColor(0xf0, 0xc0, 0x60)
TEAL   = RGBColor(0x1a, 0x8c, 0x8c)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txbox(slide, text, x, y, w, h,
          size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          italic=False, wrap=True):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf.text_frame.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tf

def rect(slide, x, y, w, h, fill_color, text="", tsize=14, tbold=False,
         tcolor=WHITE, line_color=None):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(tsize)
        run.font.bold = tbold
        run.font.color.rgb = tcolor
    return shape

def slide_header(slide, title, subtitle=""):
    rect(slide, 0, 0, 13.33, 1.0, BLUE)
    txbox(slide, title, 0.3, 0.08, 12.5, 0.75, size=28, bold=True,
          color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txbox(slide, subtitle, 0.3, 0.72, 12.5, 0.35, size=13,
              color=LGREY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE — LABELS AND FEATURES
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "Anomaly Dataset — Labels and Features",
             "5,621 labelled samples · 19 features per sample · 4 anomaly classes")

# ── LEFT: Labels ─────────────────────────────────────────────────────────────
txbox(sl, "Labels (4 classes)", 0.4, 1.1, 5.8, 0.4,
      size=15, bold=True, color=YELLOW)

labels = [
    (GREEN,  "0 — Normal",
             "Network running cleanly.\nNo stress, no impairment.\n2,063 samples"),
    (ORANGE, "1 — Scheduler Fault",
             "gNB real-time thread priority\nlowered by the OS — misses\n1ms slot deadlines.\n462 samples"),
    (RED,    "2 — Traffic Flood",
             "UE blasted with UDP traffic\nexceeding uplink capacity.\nSimulates a DDoS overload.\n936 samples"),
    (PURPLE, "3 — Channel Degradation",
             "RF impairments injected via\nthe channel broker: fading,\nnoise, burst drops.\n2,160 samples"),
]

for i, (col, title, body) in enumerate(labels):
    y = 1.55 + i * 1.35
    rect(sl, 0.4, y, 5.8, 1.28, RGBColor(0x08, 0x08, 0x1e), line_color=col)
    rect(sl, 0.4, y, 5.8, 0.38, col)
    txbox(sl, title, 0.5, y+0.04, 5.6, 0.3, size=12, bold=True, color=WHITE)
    txbox(sl, body,  0.5, y+0.42, 5.6, 0.82, size=11, color=LGREY)

# ── RIGHT: Features ───────────────────────────────────────────────────────────
txbox(sl, "Features (19 per sample)", 6.5, 1.1, 6.5, 0.4,
      size=15, bold=True, color=YELLOW)

# Hook latency group
rect(sl, 6.5, 1.55, 6.5, 0.3, PURPLE)
txbox(sl, "Hook Latency (p99) — 8 features", 6.6, 1.57, 6.3, 0.25,
      size=11, bold=True, color=WHITE)
hook_feats = [
    "fapi_ul_tti_request p99   ← key fault detector",
    "fapi_dl_tti_request p99",
    "pdcp_ul_deliver_sdu p99",
    "pdcp_ul_rx_data_pdu p99",
    "rlc_ul_rx_pdu p99",
    "rlc_ul_sdu_delivered p99",
    "rlc_dl_tx_pdu p99",
    "fapi_ul_tti_request max",
]
for i, f in enumerate(hook_feats):
    col = YELLOW if "key fault" in f else LGREY
    txbox(sl, ("  * " if "key fault" in f else "  · ") + f,
          6.5, 1.87 + i * 0.26, 6.5, 0.25, size=10, color=col)

# MAC / Radio group
rect(sl, 6.5, 4.0, 6.5, 0.3, BLUE)
txbox(sl, "MAC / Radio — 8 features", 6.6, 4.02, 6.3, 0.25,
      size=11, bold=True, color=WHITE)
mac_feats = [
    "harq_mcs_avg      — average modulation order",
    "harq_mcs_min      — worst-slot modulation",
    "harq_cons_retx    — consecutive retransmissions",
    "harq_fail_rate    — HARQ failure rate",
    "crc_sinr_avg      — average signal quality",
    "crc_harq_fail     — CRC-level HARQ failures",
    "crc_success_rate  — transmission success rate",
    "bsr_kb            — UE uplink buffer size",
]
for i, f in enumerate(mac_feats):
    txbox(sl, "  · " + f, 6.5, 4.32 + i * 0.26, 6.5, 0.25, size=10, color=LGREY)

# RLC group
rect(sl, 6.5, 6.42, 6.5, 0.3, TEAL)
txbox(sl, "RLC Layer — 3 features", 6.6, 6.44, 6.3, 0.25,
      size=11, bold=True, color=WHITE)
rlc_feats = [
    "rlc_throughput_kb  · rlc_lat_avg_us  · rlc_lat_max_us",
]
txbox(sl, "  · " + rlc_feats[0], 6.5, 6.74, 6.5, 0.25, size=10, color=LGREY)

# bottom bar
rect(sl, 0.4, 7.1, 12.5, 0.3, RGBColor(0x0a, 0x2a, 0x0a),
     "Key insight: hook_p99_us_fapi_ul is the only feature that detects scheduler faults "
     "— invisible to all standard metrics",
     tsize=11, tbold=True, tcolor=GREEN, line_color=GREEN)

prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")

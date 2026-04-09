"""
Generate a PowerPoint (.pptx) midterm presentation for the srsRAN jBPF thesis.
Structured to match the 12-week Bachelor Project Plan.
Target: 20 minutes (~17 content slides at ~1–1.5 min each).

Usage:  python3 scripts/generate_pptx.py
Output: docs/MIDTERM_PRESENTATION.pptx
"""

import os, io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── paths ──────────────────────────────────────────────────────────────────
ROOT   = os.path.dirname(os.path.abspath(__file__)) + "/.."
FIG    = ROOT + "/docs/figures"
PS     = ROOT + "/datasets/stress_anomaly/plots"
PC     = ROOT + "/datasets/channel/plots"
OUT    = ROOT + "/docs/MIDTERM_PRESENTATION.pptx"

# ── colour palette ──────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1a, 0x1a, 0x3a)   # slide background
BLUE   = RGBColor(0x1a, 0x6f, 0xaf)
GREEN  = RGBColor(0x2e, 0x8b, 0x57)
PURPLE = RGBColor(0x6a, 0x0d, 0xad)
ORANGE = RGBColor(0xe6, 0x5c, 0x00)
RED    = RGBColor(0xb2, 0x22, 0x22)
GOLD   = RGBColor(0x8b, 0x69, 0x14)
WHITE  = RGBColor(0xff, 0xff, 0xff)
LGREY  = RGBColor(0xcc, 0xcc, 0xdd)
YELLOW = RGBColor(0xf0, 0xc0, 0x60)
TEAL   = RGBColor(0x1a, 0x8c, 0x8c)

# slide dimensions: 16:9
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── helpers ────────────────────────────────────────────────────────────────

def new_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=NAVY):
    """Set slide background colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txbox(slide, text, x, y, w, h,
          size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          italic=False, wrap=True):
    """Add a text box. Returns the text frame."""
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf.text_frame.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tf

def txbox_multi(slide, lines, x, y, w, h, size=16, color=WHITE, indent=False):
    """Add a text box with multiple bullet lines."""
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf.text_frame.word_wrap = True
    first = True
    for bullet, line in lines:
        if first:
            p = tf.text_frame.paragraphs[0]
            first = False
        else:
            p = tf.text_frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if bullet:
            p.level = 1 if indent else 0
        run = p.add_run()
        run.text = ("  • " if bullet else "") + line
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tf

def rect(slide, x, y, w, h, fill_color, text="", tsize=14, tbold=False, tcolor=WHITE,
         line_color=None, radius=0):
    """Add a filled rectangle with optional label."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(tsize)
        run.font.bold = tbold
        run.font.color.rgb = tcolor
    return shape

def img(slide, path, x, y, w, h=None):
    """Add image at position. h=None → auto aspect ratio."""
    if not os.path.exists(path):
        rect(slide, x, y, w, h or 1.5, RED, f"[missing: {os.path.basename(path)}]", 10)
        return
    if h:
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w))

def divider(slide, y, color=BLUE):
    """Thin horizontal rule."""
    ln = slide.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(12.33), Inches(0.03))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()

def slide_header(slide, title, subtitle="", title_color=WHITE):
    """Standard slide header: coloured rule + title."""
    rect(slide, 0, 0, 13.33, 1.0, BLUE)
    txbox(slide, title, 0.3, 0.08, 12.5, 0.75, size=28, bold=True,
          color=title_color, align=PP_ALIGN.LEFT)
    if subtitle:
        txbox(slide, subtitle, 0.3, 0.72, 12.5, 0.35, size=13,
              color=LGREY, italic=True)

def tag_rect(slide, label, x, y, w=1.4, h=0.28, color=GREEN):
    rect(slide, x, y, w, h, color, label, tsize=11, tbold=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl, NAVY)
rect(sl, 0, 0, 13.33, 7.5, RGBColor(0x0d, 0x0d, 0x22))
rect(sl, 0, 2.5, 13.33, 0.06, BLUE)
rect(sl, 0, 4.8, 13.33, 0.06, BLUE)

txbox(sl, "Real-Time Telemetry in Open AI RAN", 0.5, 1.1, 12.3, 1.2,
      size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "using Microsoft Janus (jBPF)", 0.5, 2.15, 12.3, 0.7,
      size=26, bold=False, color=LGREY, align=PP_ALIGN.CENTER)
txbox(sl, "Midterm Progress Report", 0.5, 3.0, 12.3, 0.55,
      size=20, italic=True, color=YELLOW, align=PP_ALIGN.CENTER)
txbox(sl, "Maxim Bitiu   ·   Bachelor Thesis   ·   April 2026",
      0.5, 4.95, 12.3, 0.5, size=16, color=LGREY, align=PP_ALIGN.CENTER)
txbox(sl, "srsRAN Project  ·  jBPF / Microsoft Janus  ·  ZMQ Virtual Radio  ·  InfluxDB  ·  Grafana",
      0.5, 5.5, 12.3, 0.4, size=13, color=RGBColor(0x88, 0x88, 0xaa),
      align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — MOTIVATION
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "Why Build This?",
             "The problem with standard 5G network monitoring")

txbox_multi(sl, [
    (False, "Standard 5G monitoring gives you ~1-second averages at the MAC layer only:"),
    (True,  '"MCS dropped from 27 to 22" — but why? Radio problem? Scheduler overloaded? Traffic flood?'),
    (True,  "No way to tell from standard metrics alone"),
    (False, ""),
    (False, "eBPF (extended Berkeley Packet Filter) lets us inject probes into any running program:"),
    (True,  "Zero kernel changes — probes load and unload at runtime"),
    (True,  "Per-event granularity at every 1 ms radio slot"),
    (True,  "Multiple protocol layers simultaneously"),
    (False, ""),
    (False, "Microsoft Janus (jBPF) extends eBPF to userspace gNB software:"),
    (True,  "Hook directly into the srsRAN gNB source code at 22 function call sites"),
    (True,  "Capture microsecond-level events invisible to any standard interface"),
], 0.4, 1.1, 8.5, 5.8, size=15)

rect(sl, 9.2, 1.2, 3.8, 5.5, RGBColor(0x10, 0x10, 0x2a),
     line_color=BLUE)
txbox(sl, "Research Question", 9.3, 1.35, 3.6, 0.4, size=13, bold=True, color=YELLOW)
txbox(sl, "Can jBPF hook telemetry distinguish infrastructure faults from channel impairments in a live 5G NR system?",
      9.3, 1.8, 3.5, 2.5, size=13, color=WHITE)
txbox(sl, "Answer: Yes — and the standard metrics channel cannot do this at all.",
      9.3, 5.0, 3.5, 1.3, size=12, bold=True, color=GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SYSTEM OVERVIEW (WHAT WE BUILT)
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "What We Built",
             "A complete 5G NR telemetry pipeline — no physical radio hardware required")

img(sl, FIG + "/fig_system_architecture.png", 0.3, 1.1, 12.7, 5.7)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4a — WHY k3d DIDN'T WORK
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "First Attempt: k3d Kubernetes Deployment",
             "The official Microsoft deployment path — and why it didn't work for us")

# Intro box
rect(sl, 0.4, 1.1, 12.5, 0.75, RGBColor(0x0a, 0x20, 0x40),
     "The jrtc-apps repo (Microsoft) recommends k3d as the default deployment — "
     "lightweight Kubernetes running locally, no physical hardware required.",
     tsize=14, tbold=False, tcolor=LGREY, line_color=BLUE)

# Left: what k3d promised
txbox(sl, "What k3d Offers", 0.4, 2.05, 5.8, 0.4, size=15, bold=True, color=BLUE)
txbox_multi(sl, [
    (True, "Single Helm chart deploys gNB + jrtc + decoder + UE"),
    (True, "ZMQ mode available — no radio hardware needed"),
    (True, "Designed to run on a single laptop or workstation"),
    (True, "Containers handle all dependencies automatically"),
], 0.4, 2.45, 5.8, 2.2, size=13)

# Right: what went wrong
txbox(sl, "What Went Wrong", 6.9, 2.05, 5.9, 0.4, size=15, bold=True, color=RED)
txbox_multi(sl, [
    (True, "Helm chart values.yaml was built for enterprise O-RAN servers"),
    (True, "Pods stayed in Pending — Kubernetes couldn't schedule them"),
    (True, "Required resources our machine couldn't provide:"),
    (False, "  · Hugepages (4 Gi+) — not configured on the k3d node"),
    (False, "  · SR-IOV device plugin — Intel NIC feature not present"),
    (False, "  · DPDK CPU pinning — assumes 30+ dedicated physical cores"),
    (False, "  · Physical RU MAC/VLAN bindings in the config"),
], 6.9, 2.45, 5.9, 3.0, size=13)

# Bottom conclusion
rect(sl, 0.4, 5.8, 12.5, 0.9, RGBColor(0x30, 0x10, 0x10),
     "The chart worked as designed — it was designed for a different class of machine. "
     "We needed a different approach.",
     tsize=14, tbold=True, tcolor=YELLOW, line_color=RED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4b — BAREMETAL: WHAT WE DID
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "The Fix: Bare Metal Deployment on Ubuntu 22.04",
             "Skipped the Helm chart entirely — built and ran everything natively")

# Step boxes across the slide
steps = [
    (BLUE,   "① Build from Source",
             "Compiled srsRAN (jBPF fork), jBPF runtime, and jrtc manually — "
             "required LLVM/Clang toolchain and specific kernel headers"),
    (PURPLE, "② Fix the gNB Config",
             "Example YAML had invalid keys that silently crashed the gNB at startup. "
             "Debugged logs, identified correct keys, rewrote the config for ZMQ mode"),
    (GOLD,   "③ Fix eBPF Verifier Errors",
             "All 10 MAC codelets failed the eBPF verifier and refused to load. "
             "Fixed map definitions and removed unsupported hash clear loops"),
    (TEAL,   "④ Fix jrtc Decoder Bug",
             "DL and UL HARQ codelets shared the same schema name — the decoder "
             "silently dropped one stream. Fixed in the jrtc Go source"),
    (GREEN,  "⑤ Automate Startup",
             "Built a single launch script that starts all 8 components in correct "
             "order with dependency checks, timeouts, and clean shutdown"),
]

box_w = 2.45
for i, (col, title, body) in enumerate(steps):
    x = 0.3 + i * (box_w + 0.12)
    rect(sl, x, 1.15, box_w, 5.5, col, "", line_color=WHITE)
    txbox(sl, title, x + 0.1, 1.25, box_w - 0.2, 0.55,
          size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, body, x + 0.1, 1.95, box_w - 0.2, 4.4,
          size=12, bold=False, color=WHITE)

rect(sl, 0.4, 6.8, 12.5, 0.5, RGBColor(0x0a, 0x2a, 0x0a),
     "Result: full pipeline running — all 17 telemetry streams flowing, "
     "36-panel Grafana dashboard live, ZMQ channel broker active",
     tsize=13, tbold=True, tcolor=GREEN, line_color=GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — JBPF HOOK COVERAGE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "What We Hooked — 22 Instrumentation Points, 17 Schemas",
             "Coverage across every layer of the 5G gNB software stack")

layers = [
    # (x, y, w, h, label, sub, color)
    (0.3, 1.2, 2.5, 4.8, "MAC Scheduler", "", BLUE),
    (3.1, 1.2, 2.3, 4.8, "FAPI\n(PHY↔MAC)", "", RGBColor(0x0a, 0x5a, 0x8a)),
    (5.7, 1.2, 2.3, 4.8, "RLC / PDCP", "", PURPLE),
    (8.3, 1.2, 2.3, 4.8, "Control Plane\nRRC / NGAP", "", TEAL),
    (10.9, 1.2, 2.1, 4.8, "Performance\nMeta", "", RED),
]
for x, y, w, h, lbl, sub, col in layers:
    rect(sl, x, y, w, h, col, "", line_color=WHITE)
    txbox(sl, lbl, x+0.05, y+0.1, w-0.1, 0.5, size=13, bold=True, color=WHITE,
          align=PP_ALIGN.CENTER)

# MAC items
mac_items = ["harq_stats\n(MCS, retx,\nHARQ state)",
             "crc_stats\n(SINR, CRC\nsuccess rate)",
             "bsr_stats\n(UE buffer\noccupancy)",
             "uci_stats\n(CQI, SR,\nTiming advance)"]
for i, item in enumerate(mac_items):
    rect(sl, 0.35, 1.85 + i*0.9, 2.4, 0.82, RGBColor(0x0d, 0x3a, 0x6a), item,
         tsize=11, line_color=LGREY)

# FAPI items
fapi_items = ["fapi_dl_config\n(DL MCS, PRB,\nTBS per slot)",
              "fapi_ul_config\n(UL scheduler\ndecisions)",
              "fapi_crc_stats\n(PHY CRC,\nSNR hist)"]
for i, item in enumerate(fapi_items):
    rect(sl, 3.15, 1.85 + i*1.1, 2.2, 0.95, RGBColor(0x05, 0x30, 0x5a), item,
         tsize=11, line_color=LGREY)

# RLC/PDCP items
rlc_items = ["rlc_ul/dl_stats\n(SDU latency,\nretx bytes)",
             "pdcp_ul/dl_stats\n(per-bearer\nbyte counters)"]
for i, item in enumerate(rlc_items):
    rect(sl, 5.75, 1.85 + i*1.4, 2.2, 1.2, RGBColor(0x38, 0x07, 0x60), item,
         tsize=11, line_color=LGREY)

# Control items
ctrl_items = ["rach_stats\n(per-RACH SNR\n+ timing adv.)",
              "rrc_procedure\n(setup, reconfig,\nreestablishment)",
              "ngap_procedure\n(core network\nround-trip)"]
for i, item in enumerate(ctrl_items):
    rect(sl, 8.35, 1.85 + i*1.05, 2.2, 0.9, RGBColor(0x07, 0x42, 0x42), item,
         tsize=11, line_color=LGREY)

# Meta
rect(sl, 10.95, 1.85, 2.0, 2.3,
     RGBColor(0x5a, 0x10, 0x10),
     "jbpf_out_\nperf_list\n\np50/p90/\np95/p99\nper hook\n\n★ Novel metric\nno standard\nequivalent",
     tsize=11, line_color=YELLOW)

txbox(sl, "★  jbpf_out_perf_list is the ONLY signal sensitive to infrastructure faults — "
         "completely invisible to any standard 5G monitoring interface",
      0.3, 6.15, 12.7, 0.6, size=13, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — GRAFANA DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "36-Panel Real-Time Dashboard",
             "Live Grafana dashboard — every metric from every layer in one view")

txbox_multi(sl, [
    (False, "Dashboard sections:"),
    (True, "Custom SINR Analytics (5 panels) — variance, sliding avg, fill level"),
    (True, "PHY / MAC — Radio Link Quality (4 panels) — SINR, TX success, HARQ, RSRP"),
    (True, "MAC Scheduler (5 panels) — BSR, CQI, scheduling requests, MCS, retx"),
    (True, "FAPI (5 panels) — per-slot MCS/PRB/TBS DL vs UL, SNR, timing advance"),
    (True, "RLC / PDCP (6 panels) — DL/UL data volume, retx, SDU latency"),
    (True, "RRC / NGAP (2 panels) — control plane event timelines"),
    (True, "jBPF hook latency (3 panels) — p50/p99 per hook, invocation count"),
    (True, "Summary Statistics (6 panels) — SINR, TX%, HARQ, CQI, p50, BSR"),
    (False, ""),
    (False, "Auto-refresh: 5 s  ·  InfluxDB 1.x backend  ·  InfluxQL queries"),
], 0.4, 1.15, 5.8, 5.6, size=13)

img(sl, FIG + "/screenshot_grafana1.png", 6.1, 1.1, 7.0, 5.7)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7a — THE CHANNEL PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "The Channel Problem — Why a Broker Was Needed",
             "ZMQ virtual radio gives a perfect channel — useless for a thesis on anomaly detection")

# Left: problem statement
rect(sl, 0.3, 1.15, 6.1, 1.8, RGBColor(0x2a, 0x0a, 0x0a), line_color=RED)
txbox(sl, "The Problem", 0.45, 1.2, 3.0, 0.35, size=13, bold=True, color=RED)
txbox_multi(sl, [
    (True, "srsRAN ZMQ transport = perfect wire: SINR constant at 42 dB, zero fading, zero loss"),
    (True, "No RF variation → no channel anomalies → nothing to detect"),
    (True, "Every dataset would look identical — the thesis has no content"),
], 0.4, 1.55, 5.9, 1.3, size=12)

# srsUE built-in emulator — why it doesn't work
rect(sl, 0.3, 3.1, 6.1, 1.55, RGBColor(0x1a, 0x1a, 0x0a), line_color=ORANGE)
txbox(sl, "Attempt 1: srsUE built-in channel emulator", 0.45, 3.15, 5.8, 0.35, size=13, bold=True, color=ORANGE)
txbox_multi(sl, [
    (True, "srsUE has a --channel.model flag (AWGN, EPA, EVA, ETU)"),
    (True, "Problem: it is wired to the LTE PHY path only — silently ignored in NR mode"),
    (True, "Confirmed by reading source: NR PHY never calls the channel emulator"),
], 0.4, 3.5, 5.9, 1.05, size=12)

# C broker solution
rect(sl, 0.3, 4.75, 6.1, 1.9, RGBColor(0x0a, 0x1a, 0x0a), line_color=GREEN)
txbox(sl, "Solution: ZMQ IQ-level Broker", 0.45, 4.8, 5.8, 0.35, size=13, bold=True, color=GREEN)
txbox_multi(sl, [
    (True, "Intercepts raw IQ samples on the ZMQ socket between gNB and UE"),
    (True, "No changes to srsRAN source — broker is a transparent proxy"),
    (True, "C broker: AWGN + Rician fading (fast, low latency)"),
    (True, "Pure Rayleigh crash: |h|²→0 nulls cause RRC timeout in ~2 min"),
    (True, "Fix → Rician K=3 dB adds LoS component, prevents complete nulls"),
    (True, "Sweet spot: K=3 dB, SNR=28 dB, fd=5 Hz → stable 15+ min, 12 HARQ failures"),
], 0.4, 5.15, 5.9, 1.45, size=12)

# Right: diagram
img(sl, FIG + "/fig_channel_broker.png", 6.6, 1.15, 6.5, 5.5)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7b — FROM C BROKER TO GNU RADIO
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "From C Broker to GNU Radio — Adding Realism",
             "Flat fading is unrealistic; frequency-selective 3GPP profiles needed for the thesis")

# Left: C broker limitation
rect(sl, 0.3, 1.15, 5.9, 2.05, RGBColor(0x2a, 0x0a, 0x0a), line_color=RED)
txbox(sl, "C Broker Limitation: Flat Fading", 0.45, 1.2, 5.6, 0.35, size=13, bold=True, color=RED)
txbox_multi(sl, [
    (True, "All OFDM subcarriers fade together by the same random gain"),
    (True, "Real 5G channels are frequency-selective: different subcarriers"),
    (True, "  experience different fading (multi-path reflections at different delays)"),
    (True, "Flat fading produces unrealistic, uniform SINR — all PRBs degraded equally"),
    (True, "3GPP test profiles (EPA, EVA, ETU) require multi-tap FIR filters, not flat"),
], 0.4, 1.55, 5.7, 1.55, size=12)

# GRC broker
rect(sl, 0.3, 3.3, 5.9, 3.4, RGBColor(0x0a, 0x1a, 0x0a), line_color=GREEN)
txbox(sl, "GRC Python Broker — Built on GNU Radio", 0.45, 3.35, 5.6, 0.35, size=13, bold=True, color=GREEN)
txbox_multi(sl, [
    (True, "EPA / EVA / ETU frequency-selective fading (scipy FIR + persistent state zi)"),
    (True, "  FFT-based per-subcarrier filtering — each tap has independent delay"),
    (True, "Carrier Frequency Offset (CFO) — Doppler-like phase rotation"),
    (True, "Burst packet drops — simulate deep fade or interference events"),
    (True, "Narrowband CW interference injection"),
    (True, "Time-varying scenario scripts (10 realistic channel sequences)"),
    (True, "Live QT GUI — view IQ constellation and spectrum in real time"),
    (True, "Processing budget: DL 413 µs · UL 503 µs < 1 ms slot"),
], 0.4, 3.7, 5.7, 2.9, size=12)

# Right: GRC live GUI screenshot
txbox(sl, "Live GRC Broker GUI", 6.5, 1.15, 6.7, 0.35, size=12, bold=True, color=GREEN)
img(sl, FIG + "/screenshot_grc.png", 6.5, 1.55, 6.7, 5.1)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7c — JANUS VS STANDARD: GOAL
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "jBPF vs Standard Metrics — What Were We Trying to Prove?",
             "Week 7 evaluation — running both telemetry channels simultaneously on the same gNB")

# Left: context
rect(sl, 0.4, 1.15, 5.8, 1.5, RGBColor(0x0a, 0x20, 0x40), line_color=BLUE)
txbox(sl, "The Setup", 0.5, 1.2, 5.6, 0.4, size=14, bold=True, color=BLUE)
txbox(sl, "We ran both telemetry systems on the same gNB at the same time:\n"
          "jBPF codelets → InfluxDB 1.x   |   Standard metrics → InfluxDB 3\n"
          "57 minutes of steady-state traffic with Rician fading channel",
      0.5, 1.6, 5.6, 0.95, size=12, color=LGREY)

# Right: 3 questions
rect(sl, 6.5, 1.15, 6.6, 1.5, RGBColor(0x0a, 0x20, 0x20), line_color=TEAL)
txbox(sl, "Three Questions", 6.6, 1.2, 6.4, 0.4, size=14, bold=True, color=TEAL)
txbox(sl, "1.  Do both systems agree where they measure the same thing?\n"
          "2.  Why do the numbers differ even for the same metric?\n"
          "3.  What can jBPF see that standard metrics simply cannot?",
      6.6, 1.6, 6.3, 0.95, size=12, color=LGREY)

# Middle: 3 findings boxes
findings = [
    (GREEN,  "Finding 1 — They Agree",
             "Where both measure the same signal (SINR, MCS, CQI, BLER), "
             "values match within 1.5%. The codelets are extracting correct data."),
    (YELLOW, "Finding 2 — Differences Are Expected",
             "Throughput shows a 1.95× gap because jBPF measures at the application "
             "layer while standard metrics measure at the MAC layer — two valid "
             "measurements of different things."),
    (RED,    "Finding 3 — jBPF Sees What Standard Cannot",
             "Hook latency, per-slot FAPI data, RLC/PDCP counters, RRC/NGAP "
             "procedure timing — none of these exist in the standard interface. "
             "Infrastructure faults are invisible without them."),
]
for i, (col, title, body) in enumerate(findings):
    x = 0.4 + i * 4.35
    rect(sl, x, 2.85, 4.1, 3.5, RGBColor(0x08, 0x08, 0x1e), line_color=col)
    rect(sl, x, 2.85, 4.1, 0.45, col)
    txbox(sl, title, x+0.1, 2.88, 3.9, 0.38, size=12, bold=True, color=WHITE)
    txbox(sl, body, x+0.1, 3.38, 3.9, 2.85, size=12, color=LGREY)

rect(sl, 0.4, 6.45, 12.5, 0.55, RGBColor(0x0a, 0x2a, 0x0a),
     "Conclusion: jBPF is not a replacement for standard metrics — it is a complementary layer "
     "that adds depth, resolution, and fault visibility that standard monitoring cannot provide.",
     tsize=12, tbold=True, tcolor=GREEN, line_color=GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — JANUS VS STANDARD: CAPABILITY COMPARISON TABLE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "jBPF vs Standard Metrics — Capability Comparison",
             "What each telemetry channel can and cannot measure")

# Column headers
col_w = [5.5, 3.3, 3.3]
col_x = [0.4, 6.1, 9.6]
headers = ["Capability", "Standard Metrics", "jBPF (Janus)"]
h_colors = [BLUE, RGBColor(0x33, 0x66, 0x99), PURPLE]
for i, (hdr, col, x) in enumerate(zip(headers, h_colors, col_x)):
    rect(sl, x, 1.1, col_w[i], 0.45, col)
    txbox(sl, hdr, x+0.1, 1.13, col_w[i]-0.2, 0.38,
          size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Table rows: (capability, standard, jbpf, highlight)
rows = [
    ("Update rate",                  "~1 s aggregates",         "Configurable — down to 1 ms per slot",  False),
    ("Per-slot (1 ms) visibility",   "No",                      "Yes",                                   False),
    ("SINR / MCS / CQI / BLER",      "Yes",                     "Yes",                                   False),
    ("Hook execution latency",        "No",                      "Yes — 22 hooks (p50/p99/max)",          True),
    ("Infrastructure fault detection","No",                      "Yes — only via hook latency",           True),
    ("Per-layer byte counters",       "No",                      "Yes — RLC and PDCP separately",         False),
    ("RLC SDU delivery latency",      "No",                      "Yes",                                   False),
    ("RRC / NGAP procedure timing",   "No",                      "Yes — per-event timestamps",            False),
    ("Scheduling latency histograms", "Yes",                     "No",                                    False),
    ("PHR / Rank Indicator",          "Yes",                     "No",                                    False),
    ("Metric count",                  "~30 fields, 1 table",     "60+ fields, 17 measurements",           False),
    ("CPU overhead",                  "Negligible",              "~3.3% of one CPU core",                 False),
]

for i, (cap, std, jbpf, highlight) in enumerate(rows):
    y = 1.65 + i * 0.44
    row_bg = RGBColor(0x1a, 0x08, 0x08) if highlight else (
             RGBColor(0x10, 0x10, 0x28) if i % 2 == 0 else RGBColor(0x0a, 0x0a, 0x1e))
    border = RED if highlight else RGBColor(0x30, 0x30, 0x50)
    rect(sl, col_x[0], y, col_w[0], 0.42, row_bg, line_color=border)
    rect(sl, col_x[1], y, col_w[1], 0.42, row_bg, line_color=border)
    rect(sl, col_x[2], y, col_w[2], 0.42, row_bg, line_color=border)
    txt_col = YELLOW if highlight else WHITE
    txbox(sl, cap,  col_x[0]+0.1, y+0.04, col_w[0]-0.2, 0.36, size=11, color=txt_col, bold=highlight)
    std_col = RED if std == "No" else (GREEN if std == "Yes" else LGREY)
    txbox(sl, std,  col_x[1]+0.1, y+0.04, col_w[1]-0.2, 0.36, size=11, color=std_col, bold=False, align=PP_ALIGN.CENTER)
    jbpf_col = GREEN if jbpf.startswith("Yes") else (RED if jbpf == "No" else LGREY)
    txbox(sl, jbpf, col_x[2]+0.1, y+0.04, col_w[2]-0.2, 0.36, size=11, color=jbpf_col, bold=False)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — THROUGHPUT GAP EXPLAINED
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "The 1.95× Throughput Difference — Not a Bug",
             "Both systems are correct — they measure at different protocol layers")

img(sl, FIG + "/fig_throughput_stack.png", 0.3, 1.1, 7.0, 5.6)

txbox_multi(sl, [
    (False, "iperf3 sends 10 Mbps of UDP payload."),
    (False, "The standard srsRAN dashboard reads the MAC-layer bitrate: 19.45 Mbps."),
    (False, ""),
    (False, "The extra 9.45 Mbps comes from:"),
    (True, "Protocol headers: IP (20B) + UDP (8B) + GTP-U (16B) + PDCP (3B) + RLC (2B)"),
    (True, "MAC control elements: BSR, PHR, timing advance"),
    (True, "HARQ overhead: scheduled bytes count even for retransmissions"),
    (True, "PDCCH/reference signals: resource grant signalling"),
    (False, ""),
    (False, "Radio metrics (SINR, MCS, CQI, BLER) agree with r > 0.88."),
    (False, "The mismatch is in throughput only — expected, not an error."),
], 7.5, 1.2, 5.6, 5.4, size=13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CPU OVERHEAD
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "CPU Overhead — Measured from the Telemetry Itself",
             "jbpf_out_perf_list records hook execution time — no external profiler needed")

img(sl, FIG + "/fig_cpu_overhead.png", 0.3, 1.1, 12.7, 5.7)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — ANOMALY DATASETS (ONE SLIDE OVERVIEW)
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "Anomaly Datasets — 5 621 Labelled Samples, 19 Features",
             "35 recorded scenarios covering 4 failure types — ready for AI/ML training")

# --- 4 class boxes across the top half ----------------------------------
class_data = [
    (GREEN,  "Class 0 — Normal",
             "The network running as intended.\nClean radio, no load stress.\n\n2 063 samples\n(baseline reference)"),
    (ORANGE, "Class 1 — Scheduler Fault",
             "We lowered the priority of the gNB\nreal-time thread in the OS. The radio\nstarts missing its 1 ms deadlines.\n\n462 samples\n(software/OS fault)"),
    (RED,    "Class 2 — Traffic Flood",
             "We blasted the UE with massive\nUDP / TCP traffic — like a denial-of-\nservice attack on the radio link.\n\n936 samples\n(overload stress)"),
    (PURPLE, "Class 3 — Channel Degradation",
             "We injected RF impairments via the\nchannel broker: fading, noise, burst\ndrops — simulating poor radio coverage.\n\n2 160 samples\n(radio environment stress)"),
]
for i, (col, title, body) in enumerate(class_data):
    x = 0.3 + i * 3.27
    rect(sl, x, 1.15, 3.1, 4.0, RGBColor(0x0d, 0x0d, 0x25), line_color=col)
    rect(sl, x, 1.15, 3.1, 0.42, col)
    txbox(sl, title, x+0.08, 1.18, 3.0, 0.38, size=12, bold=True, color=WHITE)
    txbox(sl, body, x+0.1, 1.62, 2.95, 3.4, size=11, color=LGREY)

# --- Bottom: features + numbers -----------------------------------------
rect(sl, 0.3, 5.3, 12.7, 0.5, RGBColor(0x10, 0x10, 0x30), line_color=BLUE)
txbox(sl, "19 features per 1-second window  ·  7 hook latencies (p99 µs)  ·  "
          "HARQ (MCS, retx, fail rate)  ·  CRC (SINR, success rate)  ·  "
          "BSR (buffer)  ·  RLC (throughput, latency)",
      0.4, 5.35, 12.5, 0.38, size=11, color=LGREY, align=PP_ALIGN.CENTER)

# Key finding box
rect(sl, 0.3, 5.9, 12.7, 1.05, RGBColor(0x1a, 0x0a, 0x0a), "", line_color=RED)
txbox(sl, "Key finding: 14 of 23 stress scenarios produce signatures invisible to standard metrics.",
      0.4, 5.95, 12.5, 0.42, size=14, bold=True, color=YELLOW)
txbox(sl, "hook_p99_us_fapi_ul spikes to >7 000 µs (7× the 1 ms slot budget) during scheduler demotion "
          "while SINR/MCS remain normal — impossible to detect without hook latency.",
      0.4, 6.38, 12.5, 0.42, size=12, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — THE KEY RESULT: HOOK LATENCY SIGNATURE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "The Core Result — Hook Latency Identifies Infrastructure Faults",
             "A metric that exists nowhere in standard 5G monitoring")

img(sl, FIG + "/screenshot_grafana2.png", 0.3, 1.1, 8.5, 5.2)

txbox_multi(sl, [
    (False, "What happened:"),
    (True,  "gNB threads demoted from SCHED_FIFO:96 to SCHED_BATCH"),
    (True,  "OS scheduler deprioritises them — slot processing delayed"),
    (True,  "hook_p99_us_fapi_ul: 6 µs → 7 289 µs (1200× spike)"),
    (False, ""),
    (False, "What standard metrics show:"),
    (True,  "Slight MCS drop (looks like minor channel fading)"),
    (True,  "BSR increase (looks like congestion)"),
    (True,  "No indication of a software fault"),
    (False, ""),
    (False, "What hook latency shows:"),
    (True,  "Unmistakable spike — 7× the 1 ms slot budget"),
    (True,  "Pure infrastructure signature — SINR unchanged"),
    (True,  "No channel probe can see this"),
], 9.0, 1.2, 4.1, 5.4, size=12)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — IN-NETWORK CODELET (WEEK 8)
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "In-Network Pre-Processing — Week 8 Extension",
             "Computing statistics inside the eBPF program — no userspace required")

img(sl, FIG + "/screenshot_grafana1.png", 0.3, 1.1, 8.0, 5.5)

txbox_multi(sl, [
    (False, "Custom SINR codelet (C code):"),
    (True,  "Hook: mac_sched_crc_indication — fires on every UL CRC PDU"),
    (False, ""),
    (False, "Two algorithms inside the eBPF program:"),
    (True,  "Variance: Welford online algorithm"),
    (True,  "  E[X²] − E[X]² — no stored history needed"),
    (True,  "Sliding avg: 16-entry ring buffer"),
    (True,  "  Constant time update (no inner loop)"),
    (False, ""),
    (False, "Why it matters:"),
    (True,  "Anomaly detection logic runs at the gNB itself"),
    (True,  "Sub-millisecond latency — no network hop to userspace"),
    (True,  "Bandwidth saving: send avg/variance, not every raw sample"),
    (True,  "New Grafana panels: SINR variance + sliding average"),
], 8.5, 1.2, 4.6, 5.6, size=12)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — PROGRESS VS PLAN
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "Progress vs 12-Week Plan",
             "All implementation phases complete — in Phase V (writing)")

phases = [
    ("Phase I\nWks 1–2\nFoundation",    "Setup 5G sim\n+ ZMQ link",   "DONE",    GREEN),
    ("Phase II\nWks 3–5\nCorelets",     "22 hooks\n17 schemas",        "DONE",    GREEN),
    ("Phase III\nWks 6–7\nEvaluation",  "Dashboard\nComparison",       "DONE",    GREEN),
    ("Phase IV\nWks 8–9\nExtension",    "Codelet math\nAI datasets",   "DONE",    GREEN),
    ("Phase V\nWks 10–12\nConclusion",  "Validation\nThesis writing",  "IN PROG", YELLOW),
]
for i, (title, sub, status, col) in enumerate(phases):
    x = 0.4 + i * 2.55
    rect(sl, x, 1.15, 2.4, 3.0, col if status == "DONE" else RGBColor(0x2a, 0x2a, 0x0a),
         line_color=col)
    txbox(sl, title, x+0.05, 1.2, 2.3, 1.3, size=12, bold=True, color=WHITE,
          align=PP_ALIGN.CENTER)
    txbox(sl, sub, x+0.05, 2.5, 2.3, 0.8, size=11, color=LGREY,
          align=PP_ALIGN.CENTER)
    rect(sl, x+0.3, 3.85, 1.8, 0.35, col if status == "DONE" else YELLOW,
         status, tsize=12, tbold=True)

# Deliverables table
txbox(sl, "Deliverables", 0.4, 4.4, 12.5, 0.4, size=16, bold=True, color=WHITE)
deliverables = [
    ("Source code (codelets, scripts, pipeline)", True),
    ("AI-ready datasets (CSV + HDF5, 5 621 samples)", True),
    ("Setup guide and all documentation", True),
    ("Anomaly collection report (35 scenarios)", True),
    ("Janus vs Standard comparison (live data)", True),
    ("Custom codelet documentation", True),
    ("Bachelor Thesis Report", False),
    ("Final presentation", False),
]
cols = [deliverables[:4], deliverables[4:]]
for ci, col_items in enumerate(cols):
    for ri, (item, done) in enumerate(col_items):
        mark = "✓" if done else "○"
        col  = GREEN if done else RGBColor(0x88, 0x88, 0x88)
        txbox(sl, f"{mark}  {item}", 0.4 + ci*6.5, 4.85 + ri*0.42,
              6.3, 0.4, size=13, color=col)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "Summary of Contributions")

contributions = [
    (BLUE,   "1. Full 5G NR Telemetry Pipeline",
             "22 jBPF hook points · 17 schemas · 1 ms granularity · 36-panel Grafana dashboard"),
    (GOLD,   "2. ZMQ Channel Broker",
             "GRC Python broker: AWGN + Rician/Rayleigh + EPA/EVA/ETU + interference + time-varying scenarios"),
    (PURPLE, "3. Baremetal Installation",
             "Fixed 10 eBPF verifier failures + 2 jrtc-ctl bugs + YAML config keys to get the full stack running"),
    (GREEN,  "4. Two Labelled AI-Ready Datasets",
             "5 621 samples · 19 features · 4 classes · 35 scenarios · scenario-based train/test split"),
    (RED,    "5. Key Finding — Hook Latency as Fault Detector",
             "Only metric that distinguishes scheduler faults from channel impairments · 7× spike · invisible to standard"),
    (TEAL,   "6. Custom In-Network Codelet",
             "SINR variance + sliding window computed inside eBPF at 1 ms — no userspace post-processing"),
]
for i, (col, title, body) in enumerate(contributions):
    y = 1.1 + i * 0.98
    rect(sl, 0.3, y, 0.08, 0.75, col)
    txbox(sl, title, 0.5, y, 12.5, 0.38, size=15, bold=True, color=WHITE)
    txbox(sl, body,  0.5, y+0.35, 12.5, 0.45, size=12, color=LGREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — WHAT REMAINS
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "Remaining Work — Phase V")

txbox_multi(sl, [
    (False, "Thesis Writing (Week 11):"),
    (True,  "Chapter: System architecture and implementation"),
    (True,  "Chapter: Evaluation — jBPF vs standard comparison with live plots"),
    (True,  "Chapter: Anomaly detection — dataset description, key findings"),
    (True,  "Chapter: Conclusion and future work"),
    (False, ""),
    (False, "Optional extensions (if time permits):"),
    (True,  "Train a lightweight anomaly classifier on the collected dataset"),
    (True,  "Evaluate CPU overhead at multiple traffic rates (10/25/50 Mbps)"),
    (True,  "Extend custom codelet to output an anomaly score directly"),
], 0.4, 1.2, 6.5, 5.5, size=15)

rect(sl, 7.1, 1.2, 5.9, 5.5, RGBColor(0x10, 0x10, 0x2a), line_color=BLUE)
txbox(sl, "Hardware (USRP) Bonus", 7.2, 1.3, 5.7, 0.4, size=14, bold=True, color=YELLOW)
txbox_multi(sl, [
    (False, "Project plan explicitly designs for ZMQ simulation."),
    (False, "Hardware (USRP B210/X310) is optional bonus only."),
    (False, ""),
    (False, "If time permits after thesis draft:"),
    (True,  "Replace ZMQ with real USRP radio"),
    (True,  "Run codelets on live air interface"),
    (True,  "Compare simulated vs real-world hook latency"),
    (False, ""),
    (False, "Risk mitigation: do NOT attempt if it jeopardises the report deadline."),
], 7.2, 1.75, 5.6, 4.7, size=12)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl, RGBColor(0x0d, 0x0d, 0x22))
rect(sl, 0, 2.8, 13.33, 0.06, BLUE)
txbox(sl, "Thank You", 0.5, 1.3, 12.3, 1.1,
      size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "Questions?", 0.5, 2.3, 12.3, 0.6,
      size=22, color=LGREY, align=PP_ALIGN.CENTER)
txbox(sl, "github.com/MaximBitiu1/srsran-telemetry-pipeline",
      0.5, 3.1, 12.3, 0.45, size=16, color=LGREY, align=PP_ALIGN.CENTER)

stats = [
    ("22",     "jBPF hook\npoints"),
    ("17",     "Telemetry\nschemas"),
    ("45",     "Grafana\npanels"),
    ("5 621",  "Labelled\nsamples"),
    ("3.3%",   "CPU\noverhead"),
    ("7 000+", "µs fault\nspike"),
]
for i, (num, lbl) in enumerate(stats):
    x = 0.7 + i * 2.0
    txbox(sl, num, x, 3.8, 1.8, 0.6, size=26, bold=True, color=YELLOW,
          align=PP_ALIGN.CENTER)
    txbox(sl, lbl, x, 4.35, 1.8, 0.5, size=11, color=LGREY,
          align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")

"""
Standalone generator for the two jBPF vs Standard Metrics slides.

Usage:  python3 scripts/generate_jbpf_vs_standard_slides.py
Output: docs/jbpf_vs_standard_slides.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── paths ──────────────────────────────────────────────────────────────────
ROOT = os.path.dirname(os.path.abspath(__file__)) + "/.."
OUT  = ROOT + "/docs/jbpf_vs_standard_slides.pptx"

# ── colour palette ──────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1a, 0x1a, 0x3a)
BLUE   = RGBColor(0x1a, 0x6f, 0xaf)
GREEN  = RGBColor(0x2e, 0x8b, 0x57)
PURPLE = RGBColor(0x6a, 0x0d, 0xad)
RED    = RGBColor(0xb2, 0x22, 0x22)
WHITE  = RGBColor(0xff, 0xff, 0xff)
LGREY  = RGBColor(0xcc, 0xcc, 0xdd)
YELLOW = RGBColor(0xf0, 0xc0, 0x60)
TEAL   = RGBColor(0x1a, 0x8c, 0x8c)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── helpers ─────────────────────────────────────────────────────────────────

def new_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txbox(slide, text, x, y, w, h,
          size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          italic=False, wrap=True):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf.text_frame.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tf

def rect(slide, x, y, w, h, fill_color, text="", tsize=14, tbold=False,
         tcolor=WHITE, line_color=None):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(tsize)
        run.font.bold = tbold
        run.font.color.rgb = tcolor
    return shape

def slide_header(slide, title, subtitle=""):
    rect(slide, 0, 0, 13.33, 1.0, BLUE)
    txbox(slide, title, 0.3, 0.08, 12.5, 0.75, size=28, bold=True,
          color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txbox(slide, subtitle, 0.3, 0.72, 12.5, 0.35, size=13,
              color=LGREY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — GOAL: WHAT WERE WE TRYING TO PROVE?
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "jBPF vs Standard Metrics — What Were We Trying to Prove?",
             "Week 7 evaluation — running both telemetry channels simultaneously on the same gNB")

# Left: setup box
rect(sl, 0.4, 1.15, 5.8, 1.5, RGBColor(0x0a, 0x20, 0x40), line_color=BLUE)
txbox(sl, "The Setup", 0.5, 1.2, 5.6, 0.4, size=14, bold=True, color=BLUE)
txbox(sl, "We ran both telemetry systems on the same gNB at the same time:\n"
          "jBPF codelets → InfluxDB 1.x   |   Standard metrics → InfluxDB 3\n"
          "57 minutes of steady-state traffic with Rician fading channel",
      0.5, 1.6, 5.6, 0.95, size=12, color=LGREY)

# Right: 3 questions box
rect(sl, 6.5, 1.15, 6.6, 1.5, RGBColor(0x0a, 0x20, 0x20), line_color=TEAL)
txbox(sl, "Three Questions", 6.6, 1.2, 6.4, 0.4, size=14, bold=True, color=TEAL)
txbox(sl, "1.  Do both systems agree where they measure the same thing?\n"
          "2.  Why do the numbers differ even for the same metric?\n"
          "3.  What can jBPF see that standard metrics simply cannot?",
      6.6, 1.6, 6.3, 0.95, size=12, color=LGREY)

# Three findings boxes
findings = [
    (GREEN,  "Finding 1 — They Agree",
             "Where both measure the same signal (SINR, MCS, CQI, BLER), "
             "values match within 1.5%. The codelets are extracting correct data."),
    (YELLOW, "Finding 2 — Differences Are Expected",
             "Throughput shows a 1.95× gap because jBPF measures at the application "
             "layer while standard metrics measure at the MAC layer — two valid "
             "measurements of different things."),
    (RED,    "Finding 3 — jBPF Sees What Standard Cannot",
             "Hook latency, per-slot FAPI data, RLC/PDCP counters, RRC/NGAP "
             "procedure timing — none of these exist in the standard interface. "
             "Infrastructure faults are invisible without them."),
]
for i, (col, title, body) in enumerate(findings):
    x = 0.4 + i * 4.35
    rect(sl, x, 2.85, 4.1, 3.5, RGBColor(0x08, 0x08, 0x1e), line_color=col)
    rect(sl, x, 2.85, 4.1, 0.45, col)
    txbox(sl, title, x+0.1, 2.88, 3.9, 0.38, size=12, bold=True, color=WHITE)
    txbox(sl, body,  x+0.1, 3.38, 3.9, 2.85, size=12, color=LGREY)

rect(sl, 0.4, 6.45, 12.5, 0.55, RGBColor(0x0a, 0x2a, 0x0a),
     "Conclusion: jBPF is not a replacement for standard metrics — it is a complementary layer "
     "that adds depth, resolution, and fault visibility that standard monitoring cannot provide.",
     tsize=12, tbold=True, tcolor=GREEN, line_color=GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — CAPABILITY COMPARISON TABLE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl)
slide_header(sl, "jBPF vs Standard Metrics — Capability Comparison",
             "What each telemetry channel can and cannot measure")

col_w = [5.5, 3.3, 3.3]
col_x = [0.4, 6.1, 9.6]
headers = ["Capability", "Standard Metrics", "jBPF (Janus)"]
h_colors = [BLUE, RGBColor(0x33, 0x66, 0x99), PURPLE]
for i, (hdr, col, x) in enumerate(zip(headers, h_colors, col_x)):
    rect(sl, x, 1.1, col_w[i], 0.45, col)
    txbox(sl, hdr, x+0.1, 1.13, col_w[i]-0.2, 0.38,
          size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Update rate",                   "~1 s aggregates",        "Configurable — down to 1 ms per slot", False),
    ("Per-slot (1 ms) visibility",    "No",                     "Yes",                                  False),
    ("SINR / MCS / CQI / BLER",       "Yes",                    "Yes",                                  False),
    ("Hook execution latency",         "No",                     "Yes — 22 hooks (p50/p99/max)",         True),
    ("Infrastructure fault detection", "No",                     "Yes — only via hook latency",          True),
    ("Per-layer byte counters",        "No",                     "Yes — RLC and PDCP separately",        False),
    ("RLC SDU delivery latency",       "No",                     "Yes",                                  False),
    ("RRC / NGAP procedure timing",    "No",                     "Yes — per-event timestamps",           False),
    ("Scheduling latency histograms",  "Yes",                    "No",                                   False),
    ("PHR / Rank Indicator",           "Yes",                    "No",                                   False),
    ("Metric count",                   "~30 fields, 1 table",   "60+ fields, 17 measurements",          False),
    ("CPU overhead",                   "Negligible",             "~3.3% of one CPU core",                False),
]

for i, (cap, std, jbpf, highlight) in enumerate(rows):
    y = 1.65 + i * 0.44
    row_bg = RGBColor(0x1a, 0x08, 0x08) if highlight else (
             RGBColor(0x10, 0x10, 0x28) if i % 2 == 0 else RGBColor(0x0a, 0x0a, 0x1e))
    border = RED if highlight else RGBColor(0x30, 0x30, 0x50)
    rect(sl, col_x[0], y, col_w[0], 0.42, row_bg, line_color=border)
    rect(sl, col_x[1], y, col_w[1], 0.42, row_bg, line_color=border)
    rect(sl, col_x[2], y, col_w[2], 0.42, row_bg, line_color=border)
    txt_col = YELLOW if highlight else WHITE
    txbox(sl, cap,  col_x[0]+0.1, y+0.04, col_w[0]-0.2, 0.36,
          size=11, color=txt_col, bold=highlight)
    std_col = RED if std == "No" else (GREEN if std == "Yes" else LGREY)
    txbox(sl, std,  col_x[1]+0.1, y+0.04, col_w[1]-0.2, 0.36,
          size=11, color=std_col, bold=False, align=PP_ALIGN.CENTER)
    jbpf_col = GREEN if jbpf.startswith("Yes") else (RED if jbpf == "No" else LGREY)
    txbox(sl, jbpf, col_x[2]+0.1, y+0.04, col_w[2]-0.2, 0.36,
          size=11, color=jbpf_col, bold=False)


# ── save ────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
